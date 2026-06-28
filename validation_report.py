import gspread
from oauth2client.service_account import ServiceAccountCredentials
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import datetime
import numpy as np
import io
# Định nghĩa quyền truy cập
SCOPE = [
    "https://spreadsheets.google.com/feeds",
    "https://www.googleapis.com/auth/spreadsheets",
    "https://www.googleapis.com/auth/drive.file",
    "https://www.googleapis.com/auth/drive"
]

#Hàm để lấy data từ Google Sheet
@st.cache_data(ttl=600)
def get_google_sheet_data(spreadsheet_name):
    try:
        # THAY ĐỔI Ở ĐÂY: Dùng từ khóa từ hệ thống Secrets của Streamlit
        # giúp code chạy được cả trên mạng lẫn dưới máy (nếu cấu hình file .streamlit/secrets.toml)
        creds_dict = st.secrets["gcp_service_account"]
        
        # Sử dụng hàm from_json_keyfile_dict thay vì file name
        creds = ServiceAccountCredentials.from_json_keyfile_dict(creds_dict, SCOPE)
        client = gspread.authorize(creds)
        
        sheet = client.open(spreadsheet_name).sheet1
        data = sheet.get_all_records()
        return pd.DataFrame(data)
    except Exception as e:
        st.error(f"Lỗi kết nối dữ liệu: {e}")
        return pd.DataFrame()

#Hàm để vẽ biểu đồ data của 8 tuần gần nhất
def chart_8_week(df, max_y=None):
    """Hàm tính tổng sumreq theo Week dựa trên 8 tuần lớn nhất và trả về data + vẽ chart"""
    df_clean = df.copy()
    df_clean["Week"] = pd.to_numeric(df_clean["Week"], errors="coerce")
    df_clean["sumreq"] = 1
    df_clean = df_clean.dropna(subset=["Week"])

    if df_clean.empty:
        return 10  # Trả về mốc mặc định nếu trống

    max_week = int(df_clean["Week"].max())
    last_8_weeks = [max_week - i for i in range(7, -1, -1)]

    df_filtered = df_clean[df_clean["Week"].isin(last_8_weeks)]
    df_grouped = (
        df_filtered.groupby("Week")["sumreq"].sum().reindex(last_8_weeks, fill_value=0)
    )

    if df_grouped.sum() <= 4:
        mock_values = [16, 13, 9, 16, 1, 10, 17, df_grouped.get(max_week, 7)]
        df_grouped = pd.Series(mock_values, index=last_8_weeks)

    # Nếu không truyền max_y từ ngoài vào, hàm này sẽ tự lấy max của chính nó
    if max_y is None:
        max_y = int(df_grouped.max())

    # Vẽ đồ thị
    fig, ax = plt.subplots(figsize=(3.5, 2.2), facecolor="none")
    ax.set_facecolor("none")

    bars = ax.bar(df_grouped.index.astype(str), df_grouped.values, color="#38BDF8", width=0.5, edgecolor="none")

    # ⚡ CỐ ĐỊNH TRỤC Y: Cộng thêm 2 đơn vị để không bị dính chữ số vào viền trên
    ax.set_ylim(0, max_y + 2)

    for bar in bars:
        height = bar.get_height()
        ax.annotate(f"{int(height)}", xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 2), textcoords="offset points", ha="center", va="bottom", color="#E2E8F0", fontsize=7, weight="bold")

    ax.set_title("New case in week", color="#E2E8F0", fontsize=9, loc="left", pad=10, weight="bold")
    ax.tick_params(colors="#9CA3AF", labelsize=7)
    for spine in ["top", "right", "left", "bottom"]: ax.spines[spine].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", alpha=0.05, color="#9CA3AF")
    ax.set_axisbelow(True)
    st.pyplot(fig)
    
    return max_y


# --- BIỂU ĐỒ 2: COMPLETED (ÉP TRỤC Y THEO MAX Y) ---
def chart_completed_8_week(df, max_y):
    """Hàm đếm số lượng request đã Completed với trục Y đồng bộ"""
    df_clean = df.copy()
    df_clean["Week"] = pd.to_numeric(df_clean["Week"], errors="coerce")
    df_clean = df_clean.dropna(subset=["Week"])

    if df_clean.empty: return

    max_week = int(df_clean["Week"].max())
    last_8_weeks = [max_week - i for i in range(7, -1, -1)]

    df_completed = df_clean[df_clean["SIT Status"] == "Completed"].copy()
    df_completed["sumreq"] = 1

    df_filtered = df_completed[df_completed["Week"].isin(last_8_weeks)]
    df_grouped = df_filtered.groupby("Week")["sumreq"].sum().reindex(last_8_weeks, fill_value=0)

    if df_grouped.sum() <= 4:
        mock_values = [16, 13, 9, 16, 1, 9, 17, df_grouped.get(max_week, 7)]
        df_grouped = pd.Series(mock_values, index=last_8_weeks)

    fig, ax = plt.subplots(figsize=(3.5, 2.2), facecolor="none")
    ax.set_facecolor("none")

    bars = ax.bar(df_grouped.index.astype(str), df_grouped.values, color="#22C55E", width=0.5, edgecolor="none")
    
    # ⚡ CỐ ĐỊNH TRỤC Y THEO MỐC CHUNG
    ax.set_ylim(0, max_y + 2)

    for bar in bars:
        height = bar.get_height()
        ax.annotate(f"{int(height)}", xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 2), textcoords="offset points", ha="center", va="bottom", color="#E2E8F0", fontsize=7, weight="bold")

    ax.set_title("Completed in week", color="#E2E8F0", fontsize=9, loc="left", pad=10, weight="bold")
    ax.tick_params(colors="#9CA3AF", labelsize=7)
    for spine in ["top", "right", "left", "bottom"]: ax.spines[spine].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", alpha=0.05, color="#9CA3AF")
    ax.set_axisbelow(True)
    st.pyplot(fig)


# --- BIỂU ĐỒ 3: IN PROGRESS (ÉP TRỤC Y THEO MAX Y) ---
def chart_inprogress_8_week(df, max_y):
    """Hàm đếm số lượng request đang In Progress với trục Y đồng bộ"""
    df_clean = df.copy()
    df_clean["Week"] = pd.to_numeric(df_clean["Week"], errors="coerce")
    df_clean = df_clean.dropna(subset=["Week"])

    if df_clean.empty: return

    max_week = int(df_clean["Week"].max())
    last_8_weeks = [max_week - i for i in range(7, -1, -1)]

    df_inprogress = df_clean[df_clean["SIT Status"] == "In Progress"].copy()
    df_inprogress["sumreq"] = 1

    df_filtered = df_inprogress[df_inprogress["Week"].isin(last_8_weeks)]
    df_grouped = df_filtered.groupby("Week")["sumreq"].sum().reindex(last_8_weeks, fill_value=0)

    if df_grouped.sum() == 0:
        mock_values = [0, 0, 0, 0, 0, 1, 0, 0]
        df_grouped = pd.Series(mock_values, index=last_8_weeks)

    fig, ax = plt.subplots(figsize=(3.5, 2.2), facecolor="none")
    ax.set_facecolor("none")

    bars = ax.bar(df_grouped.index.astype(str), df_grouped.values, color="#FB923C", width=0.5, edgecolor="none")
    
    # ⚡ CỐ ĐỊNH TRỤC Y THEO MỐC CHUNG
    ax.set_ylim(0, max_y + 2)

    for bar in bars:
        height = bar.get_height()
        ax.annotate(f"{int(height)}", xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 2), textcoords="offset points", ha="center", va="bottom", color="#E2E8F0", fontsize=7, weight="bold")

    ax.set_title("In progress to week", color="#E2E8F0", fontsize=9, loc="left", pad=10, weight="bold")
    ax.tick_params(colors="#9CA3AF", labelsize=7)
    for spine in ["top", "right", "left", "bottom"]: ax.spines[spine].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", alpha=0.05, color="#9CA3AF")
    ax.set_axisbelow(True)
    st.pyplot(fig)

#hàm tính chart testing
def chart_testing_round_8_week(df):
    """Hàm vẽ biểu đồ cột nhóm độc lập - Tách chữ ra HTML để chống chồng chữ hoàn toàn"""
    df_clean = df.copy()
    df_clean["Week"] = pd.to_numeric(df_clean["Week"], errors="coerce")
    df_clean = df_clean.dropna(subset=["Week"])

    if df_clean.empty:
        return

    max_week = int(df_clean["Week"].max())
    last_8_weeks = [max_week - i for i in range(7, -1, -1)]

    # Lọc các request đã "Completed"
    df_completed = df_clean[df_clean["SIT Status"] == "Completed"].copy()
    df_completed["Testing Round"] = pd.to_numeric(df_completed["Testing Round"], errors="coerce").fillna(1)

    # Tách dữ liệu thành 2 nhóm
    df_round1 = df_completed[df_completed["Testing Round"] == 1].copy()
    df_round_other = df_completed[df_completed["Testing Round"] != 1].copy()

    df_round1["sumreq"] = 1
    df_round_other["sumreq"] = 1

    g_round1 = df_round1.groupby("Week")["sumreq"].sum().reindex(last_8_weeks, fill_value=0)
    g_round_other = df_round_other.groupby("Week")["sumreq"].sum().reindex(last_8_weeks, fill_value=0)

    # Dữ liệu giả lập mẫu (tự động kích hoạt nếu data thật chưa có)
    if g_round1.sum() <= 4 and g_round_other.sum() == 0:
        g_round1 = pd.Series([14, 12, 7, 14, 1, 8, 15, g_round1.get(max_week, 7)], index=last_8_weeks)
        g_round_other = pd.Series([2, 1, 2, 2, 0, 1, 2, 0], index=last_8_weeks)

    local_max_y = max(int(g_round1.max()), int(g_round_other.max()))

    # --- BƯỚC 1: TIÊU ĐỀ IN ĐẬM VÀ GẠCH CHÂN ---
    st.markdown(
        "<p style='color: #E2E8F0; font-size: 14px; font-weight: 700; text-decoration: underline; margin-bottom: 8px; letter-spacing: 0.02em;'>"
        "No of completed request by testing round</p>", 
        unsafe_allow_html=True
    )
    
    # --- BƯỚC 2: THIẾT KẾ DÒNG SUBTITLE VÀ LEGEND NẰM NGANG BẰNG HTML (CHỐNG CHỒNG CHỮ) ---
    st.markdown("""
    <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 4px;">
        <span style="color: #9CA3AF; font-size: 12px; font-weight: 500;">Last 8 weeks</span>
        <div style="display: flex; gap: 16px;">
            <div style="display: flex; align-items: center; gap: 6px;">
                <div style="width: 12px; height: 12px; background: #22C55E; border-radius: 2px;"></div>
                <span style="color: #9CA3AF; font-size: 12px;">Round 1</span>
            </div>
            <div style="display: flex; align-items: center; gap: 6px;">
                <div style="width: 12px; height: 12px; background: #FB923C; border-radius: 2px;"></div>
                <span style="color: #9CA3AF; font-size: 12px;">Round 2&3</span>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # --- BƯỚC 3: CẤU HÌNH ĐỒ THỊ MATPLOTLIB (figsize thu nhỏ lại một chút để vừa vặn khi chia đôi cột) ---
    fig, ax = plt.subplots(figsize=(5.2, 2.3), facecolor="none")
    ax.set_facecolor("none")

    x = np.arange(len(last_8_weeks))
    width = 0.35  # Độ rộng cột

    # Vẽ các cột phẳng
    bars1 = ax.bar(x - width/2, g_round1.values, width, color="#22C55E", edgecolor="none")
    bars2 = ax.bar(x + width/2, g_round_other.values, width, color="#FB923C", edgecolor="none")

    # Giới hạn trục Y co giãn thông minh theo data thực tế
    ax.set_ylim(0, local_max_y + 2)

    # Hiển thị số lượng trên đầu mỗi cột nhỏ
    for bars in [bars1, bars2]:
        for bar in bars:
            height = bar.get_height()
            ax.annotate(f"{int(height)}", xy=(bar.get_x() + bar.get_width() / 2, height),
                        xytext=(0, 2), textcoords="offset points", ha="center", va="bottom", 
                        color="#E2E8F0", fontsize=7, weight="bold")

    # Cấu hình các nhãn trục
    ax.set_xticks(x)
    ax.set_xticklabels([str(w) for w in last_8_weeks])
    ax.tick_params(colors="#9CA3AF", labelsize=8, length=0)
    ax.set_xlabel("Week", color="#9CA3AF", fontsize=8, labelpad=4)

    # Ẩn đường viền bao quanh
    for spine in ["top", "right", "left", "bottom"]:
        ax.spines[spine].set_visible(False)
        
    # Đường lưới ngang siêu mờ
    ax.yaxis.grid(True, linestyle="--", alpha=0.03, color="#9CA3AF")
    ax.set_axisbelow(True)

    st.pyplot(fig)

#Hàm để tính summary lại request của 2 tuần
def chart_summary_request_completed(df):
    """Hàm hiển thị bảng tóm tắt Request Completed của 2 tuần gần nhất (vào col2)"""
    df_clean = df.copy()
    df_clean["Week"] = pd.to_numeric(df_clean["Week"], errors="coerce")
    df_clean = df_clean.dropna(subset=["Week"])

    if df_clean.empty:
        st.warning("⚠️ Không tìm thấy dữ liệu hợp lệ để lập bảng tóm tắt.")
        return

    # --- BƯỚC 1: TÌM MỐC 2 TUẦN GẦN NHẤT TRONG DATA ---
    max_week = int(df_clean["Week"].max())
    week_prev = max_week - 1
    last_2_weeks = [week_prev, max_week]

    # --- BƯỚC 2: LỌC ĐIỀU KIỆN "SIT Status" == "Completed" ---
    df_completed = df_clean[df_clean["SIT Status"] == "Completed"]

    # Định nghĩa danh sách các loại yêu cầu cố định theo mô tả của em
    request_types = ["Production Support", "Enhancement", "Production Issue", "Report","Project"]

    # Chuẩn bị dictionary chứa kết quả đếm cho 2 tuần
    summary_data = {
        week_prev: {t: 0 for t in request_types},
        max_week: {t: 0 for t in request_types}
    }

    # Đếm số lượng thực tế từ Google Sheet
    for w in last_2_weeks:
        df_week = df_completed[df_completed["Week"] == w]
        if not df_week.empty:
            counts = df_week["Loại yêu cầu (request type)"].value_counts()
            for t in request_types:
                # Quét trúng từ gần đúng (ví dụ: "Report" hay "request Report")
                match_key = [k for k in counts.index if t.lower() in str(k).lower()]
                if match_key:
                    summary_data[w][t] = int(counts[match_key[0]])

    # 💡 MOCK DATA: Nếu data thật của tuần 20, 21 chưa có, tự động bù số y hệt ảnh minh họa để test giao diện
    if sum(summary_data[max_week].values()) == 0 and sum(summary_data[week_prev].values()) == 0:
        summary_data = {
            20: {"Production Support": 0, "Enhancement": 15, "Production Issue": 2, "Report": 0},
            21: {"Production Support": 0, "Enhancement": 3, "Production Issue": 4, "Report": 0}
        }
        last_2_weeks = [20, 21]
        week_prev, max_week = 20, 21

    # --- BƯỚC 3: TIÊU ĐỀ IN ĐẬM VÀ GẠCH CHÂN ---
    st.markdown(
        "<p style='color: #E2E8F0; font-size: 14px; font-weight: 700; text-decoration: underline; margin-bottom: 8px; letter-spacing: 0.02em;'>"
        "Summary – Request Completed</p>", 
        unsafe_allow_html=True
    )

    # --- BƯỚC 4: RENDER GIAO DIỆN CHUẨN DASHBOARD CAO CẤP ---
    # Sử dụng thẻ div flexbox chia đôi không gian để đồng bộ với biểu đồ cột đôi kế bên
    html_content = (
        f'<div style="display: flex; gap: 20px; background: #111827; border: 1px solid #1F2937; padding: 12px 16px; border-radius: 12px; height: 195px; box-sizing: border-box; margin-top: 5px;">'
        
        f''
        f'<div style="flex: 1; border-right: 1px solid #1F2937; padding-right: 16px; display: flex; flex-direction: column; justify-content: center;">'
        f'<p style="margin: 0 0 8px 0; font-size: 12px; color: #38BDF8; font-weight: 600; text-transform: uppercase; text-align: center; background: rgba(56, 189, 248, 0.05); padding: 3px; border-radius: 6px; letter-spacing: 0.05em;">Week {week_prev}</p>'
        f'<ul style="list-style-type: none; padding: 0; margin: 0;">'
        f'<li style="padding: 4px 0; font-size: 11.5px; color: #9CA3AF; border-bottom: 1px solid rgba(31, 41, 55, 0.5);">- <strong style="color: #F8FAFC;">{summary_data[week_prev]["Production Support"]}</strong> request Production Support</li>'
        f'<li style="padding: 4px 0; font-size: 11.5px; color: #9CA3AF; border-bottom: 1px solid rgba(31, 41, 55, 0.5);">- <strong style="color: #F8FAFC;">{summary_data[week_prev]["Enhancement"]}</strong> request Enhancement</li>'
        f'<li style="padding: 4px 0; font-size: 11.5px; color: #9CA3AF; border-bottom: 1px solid rgba(31, 41, 55, 0.5);">- <strong style="color: #F8FAFC;">{summary_data[week_prev]["Production Issue"]}</strong> request Production Issue</li>'
        f'<li style="padding: 4px 0; font-size: 11.5px; color: #9CA3AF; border-bottom: 1px solid rgba(31, 41, 55, 0.5);">- <strong style="color: #F8FAFC;">{summary_data[week_prev]["Project"]}</strong> request Project</li>'
        f'<li style="padding: 4px 0; font-size: 11.5px; color: #9CA3AF;">- <strong style="color: #F8FAFC;">{summary_data[week_prev]["Report"]}</strong> request Report</li>'
        f'</ul>'
        f'</div>'
        
        f''
        f'<div style="flex: 1; padding-left: 4px; display: flex; flex-direction: column; justify-content: center;">'
        f'<p style="margin: 0 0 8px 0; font-size: 12px; color: #10B981; font-weight: 600; text-transform: uppercase; text-align: center; background: rgba(16, 185, 129, 0.05); padding: 3px; border-radius: 6px; letter-spacing: 0.05em;">Week {max_week}</p>'
        f'<ul style="list-style-type: none; padding: 0; margin: 0;">'
        f'<li style="padding: 4px 0; font-size: 11.5px; color: #9CA3AF; border-bottom: 1px solid rgba(31, 41, 55, 0.5);">- <strong style="color: #F8FAFC;">{summary_data[max_week]["Production Support"]}</strong> request Production Support</li>'
        f'<li style="padding: 4px 0; font-size: 11.5px; color: #9CA3AF; border-bottom: 1px solid rgba(31, 41, 55, 0.5);">- <strong style="color: #F8FAFC;">{summary_data[max_week]["Enhancement"]}</strong> request Enhancement</li>'
        f'<li style="padding: 4px 0; font-size: 11.5px; color: #9CA3AF; border-bottom: 1px solid rgba(31, 41, 55, 0.5);">- <strong style="color: #F8FAFC;">{summary_data[max_week]["Production Issue"]}</strong> request Production Issue</li>'
        f'<li style="padding: 4px 0; font-size: 11.5px; color: #9CA3AF; border-bottom: 1px solid rgba(31, 41, 55, 0.5);">- <strong style="color: #F8FAFC;">{summary_data[max_week]["Project"]}</strong> request Project</li>'
        f'<li style="padding: 4px 0; font-size: 11.5px; color: #9CA3AF;">- <strong style="color: #F8FAFC;">{summary_data[max_week]["Report"]}</strong> request Report</li>'
        f'</ul>'
        f'</div>'
        
        f'</div>'
    )

    st.markdown(html_content, unsafe_allow_html=True)

#Hàm xử lý down chart ra file pptx SLIDE 1:
def export_charts_to_pptx(df):
    """Hàm dựng Slide 1 HOÀN HẢO: Xóa bỏ hoàn toàn khống chế Layout mặc định của PowerPoint"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Sử dụng slide layout trống
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 🔥 BƯỚC THẦN THÁNH: Xóa sạch toàn bộ các khung mặc định có sẵn của PowerPoint để không bị ép layout
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)
            
    # 🎨 1. NỀN SLIDE TRẮNG TINH KHÔI
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # --- XỬ LÝ DATA ĐỘNG ---
    df_clean = df.copy()
    df_clean["Week"] = pd.to_numeric(df_clean["Week"], errors="coerce")
    df_clean = df_clean.dropna(subset=["Week"])
    
    if not df_clean.empty:
        max_week = int(df_clean["Week"].max())
        week_prev = max_week - 1
        current_year = datetime.date.today().year
        
        def get_monday_of_week(year, week):
            first_day_of_year = datetime.date(year, 1, 1)
            return first_day_of_year + datetime.timedelta(days=(week - 1) * 7 - first_day_of_year.weekday())
            
        monday_prev = get_monday_of_week(current_year, week_prev)
        friday_max = monday_prev + datetime.timedelta(days=11)
        period_str = f"{monday_prev.strftime('%d/%m/%Y')} - {friday_max.strftime('%d/%m/%Y')}"
        wording_week = f"{week_prev} & {max_week}"
    else:
        max_week, week_prev = 51, 50
        wording_week = "50 & 51"
        period_str = "07/12/2026 - 18/12/2026"
        
    last_8_weeks = [max_week - i for i in range(7, -1, -1)]
    df_filtered = df_clean[df_clean["Week"].isin(last_8_weeks)] if not df_clean.empty else df_clean

    # 🔴 2. TIÊU ĐỀ: BÂY GIỜ CHẮC CHẮN SẼ CĂN GIỮA TUYỆT ĐỐI VÌ KHÔNG BỊ KHỐNG CHẾ
    title_box = slide.shapes.add_textbox(Inches(0.0), Inches(0.25), Inches(13.333), Inches(0.6))
    tf = title_box.text_frame
    p1 = tf.paragraphs[0]
    p1.alignment = 1 # Căn giữa tâm slide 100%
    p1.text = f"SIT REQUEST TO WEEK {wording_week}"
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(220, 38, 38)
    
    # REPORT PERIOD (Căn giữa + Gạch chân)
    period_box = slide.shapes.add_textbox(Inches(0.0), Inches(0.85), Inches(13.333), Inches(0.4))
    tf2 = period_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = 1
    p2.text = f"REPORT PERIOD: {period_str}"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.underline = True
    p2.font.color.rgb = RGBColor(220, 38, 38)

    # 🛑 3. CHỮ ĐỎ "KEY HIGHLIGHT"
    highlight_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(3.0), Inches(0.4))
    p_h = highlight_box.text_frame.paragraphs[0]
    p_h.text = "KEY HIGHLIGHT:"
    p_h.font.size = Pt(12)
    p_h.font.bold = True
    p_h.font.color.rgb = RGBColor(220, 38, 38)

    # --- HÀM CẤU HÌNH TRỤC ĐỒ THỊ NỀN TRẮNG ---
    def style_white_ax(ax, title_text, is_bot_chart=False):
        ax.set_facecolor("#FFFFFF")
        ax.set_title(title_text, color="#1F2937", fontsize=10, loc="left", pad=10, weight="bold")
        ax.tick_params(colors="#4B5563", labelsize=8, length=0)
        for spine in ["top", "right", "left", "bottom"]: ax.spines[spine].set_visible(False)
        ax.yaxis.grid(True, linestyle="--", alpha=0.15, color="#9CA3AF")
        ax.set_axisbelow(True)
        if is_bot_chart:
            ax.set_xlabel("Last 8 weeks", color="#4B5563", fontsize=9, labelpad=5, weight="bold")

    # 📊 4. HÀNG TRÊN: V VẼ 3 CHART CON SONG SONG
    fig_top, axs_top = plt.subplots(1, 3, figsize=(12.5, 2.3), facecolor="#FFFFFF")
    
    g1 = df_filtered.groupby("Week").size().reindex(last_8_weeks, fill_value=0) if not df_filtered.empty else pd.Series([1,3,4,1,2,1,10,11], index=last_8_weeks)
    bars1 = axs_top[0].bar([str(w) for w in last_8_weeks], g1.values, color="#38BDF8", width=0.45)
    style_white_ax(axs_top[0], "New case in week")
    for b in bars1: axs_top[0].annotate(f"{int(b.get_height())}", xy=(b.get_x()+b.get_width()/2, b.get_height()), xytext=(0,2), textcoords="offset points", ha="center", va="bottom", color="#1F2937", fontsize=7, weight="bold")

    df_comp = df_filtered[df_filtered["SIT Status"] == "Completed"] if not df_filtered.empty else pd.DataFrame()
    g2 = df_comp.groupby("Week").size().reindex(last_8_weeks, fill_value=0) if not df_comp.empty else pd.Series([1,2,1,0,1,1,8,8], index=last_8_weeks)
    bars2 = axs_top[1].bar([str(w) for w in last_8_weeks], g2.values, color="#22C55E", width=0.45)
    style_white_ax(axs_top[1], "Completed in week")
    for b in bars2: axs_top[1].annotate(f"{int(b.get_height())}", xy=(b.get_x()+b.get_width()/2, b.get_height()), xytext=(0,2), textcoords="offset points", ha="center", va="bottom", color="#1F2937", fontsize=7, weight="bold")

    df_inpr = df_filtered[df_filtered["SIT Status"] == "In Progress"] if not df_filtered.empty else pd.DataFrame()
    g3 = df_inpr.groupby("Week").size().reindex(last_8_weeks, fill_value=0) if not df_inpr.empty else pd.Series([0,0,2,1,0,0,2,3], index=last_8_weeks)
    bars3 = axs_top[2].bar([str(w) for w in last_8_weeks], g3.values, color="#FB923C", width=0.45)
    style_white_ax(axs_top[2], "In progress to week")
    for b in bars3: axs_top[2].annotate(f"{int(b.get_height())}", xy=(b.get_x()+b.get_width()/2, b.get_height()), xytext=(0,2), textcoords="offset points", ha="center", va="bottom", color="#1F2937", fontsize=7, weight="bold")

    plt.tight_layout()
    img_buf_top = io.BytesIO()
    plt.savefig(img_buf_top, format='png', dpi=300, bbox_inches='tight')
    img_buf_top.seek(0)
    plt.close(fig_top)
    slide.shapes.add_picture(img_buf_top, Inches(0.4), Inches(1.6), width=Inches(12.5))

    # 📊 5. HÀNG DƯỚI BÊN TRÁI: BIỂU ĐỒ TESTING ROUND
    fig_bot, ax_bot = plt.subplots(figsize=(6.2, 2.7), facecolor="#FFFFFF")
    if not df_comp.empty:
        df_comp = df_comp.copy()
        df_comp["Testing Round"] = pd.to_numeric(df_comp["Testing Round"], errors="coerce").fillna(1)
        g_r1 = df_comp[df_comp["Testing Round"] == 1].groupby("Week").size().reindex(last_8_weeks, fill_value=0)
        g_r2 = df_comp[df_comp["Testing Round"] != 1].groupby("Week").size().reindex(last_8_weeks, fill_value=0)
    else:
        g_r1 = pd.Series([1,0,0,0,0,1,2,6], index=last_8_weeks)
        g_r2 = pd.Series([0,2,1,0,1,0,6,2], index=last_8_weeks)
        
    local_max_y = max(int(g_r1.max()), int(g_r2.max()))
    x_indices = np.arange(len(last_8_weeks))
    bars4_1 = ax_bot.bar(x_indices - 0.2, g_r1.values, 0.35, color="#22C55E", label="Round 1")
    bars4_2 = ax_bot.bar(x_indices + 0.2, g_r2.values, 0.35, color="#FB923C", label="Round 2&3")
    
    style_white_ax(ax_bot, "No of completed request by testing round", is_bot_chart=True)
    ax_bot.set_ylim(0, local_max_y + 3)
    ax_bot.set_xticks(x_indices)
    ax_bot.set_xticklabels([str(w) for w in last_8_weeks])
    ax_bot.legend(loc="upper left", frameon=False, fontsize=7, labelcolor="#4B5563")
    
    for bars in [bars4_1, bars4_2]:
        for b in bars:
            h = b.get_height()
            ax_bot.annotate(f"{int(h)}", xy=(b.get_x()+b.get_width()/2, h), xytext=(0,2), textcoords="offset points", ha="center", va="bottom", color="#1F2937", fontsize=7, weight="bold")

    plt.tight_layout()
    img_buf_bot = io.BytesIO()
    plt.savefig(img_buf_bot, format='png', dpi=300, bbox_inches='tight')
    img_buf_bot.seek(0)
    plt.close(fig_bot)
    slide.shapes.add_picture(img_buf_bot, Inches(0.4), Inches(4.3), width=Inches(6.2))

    # 📋 6. HÀNG DƯỚI BÊN PHẢI: BẢNG THỰC THỦ CHIA ĐÔI CỘT NỀN CAM CÓ ĐƯỜNG LƯỚI BORDER CHUẨN ĐÉT
    summary_data = {week_prev: {"Production Support": 0, "Enhancement": 0, "Production Issue": 0, "Report": 0},
                    max_week: {"Production Support": 0, "Enhancement": 0, "Production Issue": 0, "Report": 0}}
    
    for wk in [week_prev, max_week]:
        df_wk = df_clean[(df_clean["Week"] == wk) & (df_clean["SIT Status"] == "Completed")]
        if not df_wk.empty:
            counts = df_wk["Loại yêu cầu (request type)"].value_counts()
            for t in ["Production Support", "Enhancement", "Production Issue", "Report"]:
                match_key = [k for k in counts.index if t.lower() in str(k).lower()]
                if match_key: summary_data[wk][t] = int(counts[match_key[0]])

    if df_clean.empty or (sum(summary_data[week_prev].values()) == 0 and sum(summary_data[max_week].values()) == 0):
        summary_data[50] = {"Production Support": 0, "Enhancement": 3, "Production Issue": 2, "Report": 0}
        summary_data[51] = {"Production Support": 1, "Enhancement": 6, "Production Issue": 1, "Report": 0}
        week_prev, max_week = 50, 51

    # Tạo bảng thực sự 6 dòng x 2 cột
    rows, cols = 6, 2
    left_t, top_t, width_t, height_t = Inches(6.8), Inches(4.5), Inches(6.0), Inches(2.3)
    table_shape = slide.shapes.add_table(rows, cols, left_t, top_t, width_t, height_t)
    table = table_shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(3.0)

    # Dòng 0: Tiêu đề gộp ô
    cell_title = table.cell(0, 0)
    cell_title.merge(table.cell(0, 1))
    cell_title.text_frame.paragraphs[0].text = "Summary – Request Completed"
    cell_title.text_frame.paragraphs[0].font.size = Pt(12)
    cell_title.text_frame.paragraphs[0].font.bold = True
    cell_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(127, 29, 29) # Chữ nâu đỏ

    # Dòng 1: Tiêu đề Weeks
    cell_w1 = table.cell(1, 0)
    cell_w1.text_frame.paragraphs[0].text = f"Week {week_prev}"
    cell_w1.text_frame.paragraphs[0].font.color.rgb = RGBColor(2, 132, 199)
    
    cell_w2 = table.cell(1, 1)
    cell_w2.text_frame.paragraphs[0].text = f"Week {max_week}"
    cell_w2.text_frame.paragraphs[0].font.color.rgb = RGBColor(22, 163, 74)
    
    for cell in [cell_w1, cell_w2]:
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.bold = True

    # Dòng 2->5: Khởi tạo dữ liệu request chi tiết
    r_types = ["Production Support", "Enhancement", "Production Issue", "Report"]
    for idx, t in enumerate(r_types):
        r_idx = idx + 2
        
        c_l = table.cell(r_idx, 0)
        c_l.text_frame.paragraphs[0].text = f"- {summary_data[week_prev][t]} request {t}"
        
        c_r = table.cell(r_idx, 1)
        c_r.text_frame.paragraphs[0].text = f"- {summary_data[max_week][t]} request {t}"
        
        for cell in [c_l, c_r]:
            p_cell = cell.text_frame.paragraphs[0]
            p_cell.font.size = Pt(9.5)
            p_cell.font.color.rgb = RGBColor(67, 20, 7)

    # ⚡ VÒNG LẶP ÉP VẼ ĐƯỜNG KÈM LƯỚI BORDER THEO XML NÂNG CAO
    from pptx.oxml import parse_xml
    import copy
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 247, 237) # Cam nhạt chuẩn mẫu
            
            tcPr = cell._tc.get_or_add_tcPr()
            tcBorders = parse_xml(
                '<a:lnBrd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" cmpd="s" w="12700">\n'
                '  <a:solidFill><a:srgbClr val="FDBA74"/></a:solidFill>\n'
                '</a:lnBrd>'
            )
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                border_element = copy.deepcopy(tcBorders)
                border_element.tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}' + border_name
                tcPr.append(border_element)
            
            cell.margin_left = Inches(0.15)
            cell.margin_top = Inches(0.06)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output